import os
import shutil
from typing import Dict, Optional, Tuple

import fitz  # PyMuPDF
import docx
import pandas as pd  # Import pandas to read Excel and CSV files
from pptx import Presentation  # Import Presentation for PPT files
from ebooklib import epub
import ebooklib
from mobi import extract
from bs4 import BeautifulSoup

_file_cache: Dict[Tuple[str, float], Optional[str]] = {}


def read_text_file(file_path):
    """Read text content from a text file."""
    max_chars = 3000  # Limit processing time
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
            text = file.read(max_chars)
        return text
    except Exception as e:
        print(f"Error reading text file {file_path}: {e}")
        return None


def read_docx_file(file_path):
    """Read text content from a .docx or .doc file."""
    try:
        doc = docx.Document(file_path)
        full_text = [para.text for para in doc.paragraphs]
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error reading DOCX file {file_path}: {e}")
        return None


def read_pdf_file(file_path):
    """Read text content from a PDF file."""
    try:
        doc = fitz.open(file_path)
        # Read only the first few pages to speed up processing
        num_pages_to_read = 3  # Adjust as needed
        full_text = []
        for page_num in range(min(num_pages_to_read, len(doc))):
            page = doc.load_page(page_num)
            full_text.append(page.get_text())
        pdf_content = "\n".join(full_text)
        return pdf_content
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
        return None


def read_spreadsheet_file(file_path):
    """Read text content from an Excel or CSV file."""
    try:
        if file_path.lower().endswith(".csv"):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        text = df.to_string()
        return text
    except Exception as e:
        print(f"Error reading spreadsheet file {file_path}: {e}")
        return None


def read_ppt_file(file_path):
    """Read text content from a PowerPoint file."""
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error reading PowerPoint file {file_path}: {e}")
        return None


def read_epub_file(file_path):
    """Read text content from an EPUB file."""
    try:
        book = epub.read_epub(file_path)
        text_parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            text_parts.append(item.get_content().decode("utf-8", "ignore"))
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error reading EPUB file {file_path}: {e}")
        return None


def read_mobi_file(file_path):
    """Read text content from a MOBI/AZW file."""
    try:
        tempdir, extracted = extract.extract(file_path)
        ext = os.path.splitext(extracted)[1].lower()
        if ext == ".epub":
            result = read_epub_file(extracted)
        elif ext == ".html":
            with open(extracted, "r", encoding="utf-8", errors="ignore") as file:
                html_content = file.read()
            result = BeautifulSoup(html_content, "html.parser").get_text()
        elif ext == ".pdf":
            result = read_pdf_file(extracted)
        else:
            result = None
        shutil.rmtree(tempdir, ignore_errors=True)
        return result
    except Exception as e:
        print(f"Error reading MOBI file {file_path}: {e}")
        return None


def read_file_data(file_path):
    """Read content from a file based on its extension with caching."""
    key = (file_path, os.path.getmtime(file_path))
    if key in _file_cache:
        return _file_cache[key]

    ext = os.path.splitext(file_path.lower())[1]
    if ext in [".txt", ".md"]:
        result = read_text_file(file_path)
    elif ext in [".docx", ".doc"]:
        result = read_docx_file(file_path)
    elif ext == ".pdf":
        result = read_pdf_file(file_path)
    elif ext in [".xls", ".xlsx", ".csv"]:
        result = read_spreadsheet_file(file_path)
    elif ext in [".ppt", ".pptx"]:
        result = read_ppt_file(file_path)
    elif ext == ".epub":
        result = read_epub_file(file_path)
    elif ext in [".mobi", ".azw", ".azw3"]:
        result = read_mobi_file(file_path)
    else:
        result = None

    _file_cache[key] = result
    return result


def display_directory_tree(path):
    """Display a directory tree similar to the ``tree`` command.

    The output includes the full path for each item.
    """

    def tree(dir_path, prefix=""):
        contents = sorted([c for c in os.listdir(dir_path) if not c.startswith(".")])
        pointers = ["├── "] * (len(contents) - 1) + ["└── "] if contents else []
        for pointer, name in zip(pointers, contents):
            full_path = os.path.join(dir_path, name)
            print(prefix + pointer + name)
            if os.path.isdir(full_path):
                extension = "│   " if pointer == "├── " else "    "
                tree(full_path, prefix + extension)

    if os.path.isdir(path):
        print(os.path.abspath(path))
        tree(path)
    else:
        print(os.path.abspath(path))


def collect_file_paths(base_path):
    """Collect file paths from the given directory or single file.

    Hidden files are excluded from the result.
    """
    if os.path.isfile(base_path):
        return [base_path]
    else:
        file_paths = []
        for root, _, files in os.walk(base_path):
            for file in files:
                if not file.startswith("."):  # Exclude hidden files
                    file_paths.append(os.path.join(root, file))
        return file_paths


def separate_files_by_type(file_paths):
    """Separate files into images and text files based on their extensions."""
    image_extensions = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff")
    text_extensions = (
        ".txt",
        ".docx",
        ".doc",
        ".pdf",
        ".md",
        ".xls",
        ".xlsx",
        ".ppt",
        ".pptx",
        ".csv",
        ".epub",
        ".mobi",
        ".azw",
        ".azw3",
    )
    image_files = [
        fp for fp in file_paths if os.path.splitext(fp.lower())[1] in image_extensions
    ]
    text_files = [
        fp for fp in file_paths if os.path.splitext(fp.lower())[1] in text_extensions
    ]

    return image_files, text_files  # Return only two values


def clear_file_cache():
    """Clear the internal file content cache."""
    _file_cache.clear()
